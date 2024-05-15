import asyncio
import os
from pathlib import Path
from typing import Generator
from PIL import Image
import shutil
import pptx
import pptx.slide
import fitz

from src.util import Util


CWD = os.getcwd()


class File:
    Input = f"{CWD}/@input"

    Backup = f"{CWD}/backup"
    Output = f"{CWD}/_output"
    Origin = f"{Output}/original"
    Resized = f"{Output}/resized"
    ResizedPL = f"{Output}/resizedPL"

    def __init__(self, util: Util) -> None:
        self.util = util
        self.image = ImageGenerator()
        self.parser = Parser()

    def _make_dir(self, path: str):
        p = Path(path)
        if not p.exists():
            p.mkdir()

    def _delete_dir(self, dir: str):
        shutil.rmtree(dir, ignore_errors=True)
        if not Path(dir).exists():
            Path(dir).mkdir()

    def init(self):
        for v in [
            File.Input,
            File.Output,
            Parser.Pptx,
            Parser.Pdfs,
            File.Origin,
            File.Resized,
            File.ResizedPL,
            File.Backup,
        ]:
            self._make_dir(v)

    async def generate_images(
        self, li: Generator[Path, None, None], folder_name: str = ""
    ):
        resized = ""
        pl = ""

        for i, path in enumerate(li):
            if path.is_file():
                img, size = self.image.get_image_with_size(f"{path.absolute()}")

                if "." not in path.stem:
                    resized = (
                        path.name.split(".")[0].replace(" ", "")
                        + "_RESIZED"
                        + path.suffix
                    )
                    pl = path.name.split(".")[0].replace(" ", "") + "_PL" + path.suffix
                else:
                    resized: str = (
                        path.stem.replace(".", "_").replace(" ", "")
                        + "_RESIZED"
                        + path.suffix
                    )
                    pl: str = (
                        path.stem.replace(".", "_").replace(" ", "")
                        + "_PL"
                        + path.suffix
                    )

                if len(resized) == 0:
                    raise Exception("fileN isn't captured")

                async with asyncio.TaskGroup() as tg:
                    tg.create_task(
                        self.image.save_image(
                            image=img,
                            location=File.Origin,
                            fileName=path.name.replace(" ", "")
                            if not folder_name
                            else folder_name + "_" + path.name.replace(" ", ""),
                            index=i + 1,
                        )
                    )
                    tg.create_task(
                        self.image.save_image(
                            image=self.image.resize_image(
                                image=img, size=size, isPLSize=False
                            ),
                            location=File.Resized,
                            fileName=resized
                            if not folder_name
                            else folder_name + "_" + resized,
                            index=i + 1,
                        )
                    )
                    tg.create_task(
                        self.image.save_image(
                            image=self.image.resize_image(
                                image=img, size=size, isPLSize=True
                            ),
                            location=File.ResizedPL,
                            fileName=pl if not folder_name else folder_name + "_" + pl,
                            index=i + 1,
                        )
                    )
            else:
                nested = Path.iterdir(Path(path.absolute()))
                await self.generate_images(nested, folder_name=path.name)

    async def move_images(
        self,
    ):
        time_tz = self.util.get_time_tz()
        dest = Path(f"{File.Output}/{time_tz}")
        if not dest.exists():
            dest.mkdir()
        else:
            paths = [
                f"{dest}/origin",
                f"{dest}/resized",
                f"{dest}/resizedPLL",
                f"{dest}",
            ]
            for p in paths:
                self._delete_dir(p)

        async def moving(dir: str):
            shutil.move(dir, dest)

        async with asyncio.TaskGroup() as tg:
            tg.create_task(moving(File.Origin))
            tg.create_task(moving(File.Resized))
            tg.create_task(moving(File.ResizedPL))

    async def main(self):
        self.init()

        if len(os.listdir(f"{os.getcwd()}/@input")) == 0:
            raise Exception("Empty Inputs!")

        print("Parsing Images [1/3]")
        self.parser.parsing_pptx()
        self.parser.parsing_pdf()
        li = Path.iterdir(Path(File.Input))
        await self.generate_images(li)
        await self.move_images()
        self._delete_dir(File.Backup)
        shutil.move(f"{File.Input}", f"{File.Backup}")
        self._delete_dir(File.Input)


class ImageGenerator:
    Limit = 300
    PL = 60

    def get_image_with_size(self, path: str) -> tuple[Image.Image, tuple[int, int]]:
        image = Image.open(f"{path}")
        return image, image.size

    def resize_image(
        self, image: Image.Image, size: tuple[int, int], isPLSize: bool
    ) -> Image.Image:
        def is_height_bigger(size: tuple[int, int]) -> bool:
            if size[0] < size[1]:
                return True
            return False

        resize = ImageGenerator.Limit if not isPLSize else ImageGenerator.PL
        if is_height_bigger(size):
            percent_width = resize / float(image.size[0])
            new_height = int(float(image.size[1]) * float(percent_width))
            resized_image = image.resize((resize, new_height))
            return resized_image
        else:
            percent_height = resize / float(image.size[1])
            new_width = int(float(image.size[0]) * float(percent_height))
            resized_image = image.resize((new_width, resize))
            return resized_image

    async def save_image(
        self, image: Image.Image, location: str, fileName: str, index: int
    ):
        image.save(
            f"{location}/{fileName if '@' in fileName else f"{fileName.split('.')[0]}@{index}.{fileName.split('.')[-1]}"}"
        )


class Parser:
    Pdfs = f"{CWD}/pdfs"
    Pptx = f"{CWD}/pptxs"
    To = f"{CWD}/@input"

    def parsing_pdf(self):
        pdfs = Path.iterdir(Path(Parser.Pdfs))
        for pdf in pdfs:
            doc: fitz.Document = fitz.open(pdf.absolute())
            for i in range(len(doc)):
                for j, img in enumerate(doc.get_page_images(i)):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)
                    filename = f"{Parser.To}/{pdf.stem}@{j + 1}.png"
                    pix.save(filename=filename)

    def parsing_pptx(self):
        li = Path.iterdir(Path(Parser.Pptx))
        for f in li:
            prs = pptx.Presentation(f)
            name = f.stem

            i = 1
            for slide in prs.slides:
                shapes: pptx.slide.SlideShapes = slide.shapes
                for shape in shapes:
                    if hasattr(shape, "image"):
                        image: pptx.ImagePart = shape.image  # type: ignore
                        image_bytes = image.blob

                        path = Path(Parser.To)
                        if not path.exists():
                            path.mkdir()

                        path = f"{path}/{name}@{i}.jpg"
                        i += 1
                        if image_bytes is not None:
                            with open(
                                path,
                                "wb",
                            ) as f:
                                f.write(image_bytes)
                        else:
                            raise Exception("image buf is empty")
