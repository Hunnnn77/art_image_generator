from pathlib import Path
from typing import Generator
from PIL import Image
import pptx
import pptx.slide
import fitz

from src.util import CWD


class ImageGenerator:
    Limit = 300
    PL = 60

    def get_image_with_size(self, path: str) -> tuple[Image.Image, tuple[int, int]]:
        image = Image.open(f"{path}")
        return image, image.size

    def resize_image(
        self, image: Image.Image, size: tuple[int, int], isPLSize: bool
    ) -> Image.Image:
        def is_height_bigger(size: tuple[int, int]) -> bool:
            if size[0] < size[1]:
                return True
            return False

        resize = ImageGenerator.Limit if not isPLSize else ImageGenerator.PL
        if is_height_bigger(size):
            percent_width = resize / float(image.size[0])
            new_height = int(float(image.size[1]) * float(percent_width))
            resized_image = image.resize((resize, new_height))
            return resized_image
        else:
            percent_height = resize / float(image.size[1])
            new_width = int(float(image.size[0]) * float(percent_height))
            resized_image = image.resize((new_width, resize))
            return resized_image

    async def save_image(
        self, image: Image.Image, location: str, fileName: str, index: int
    ):
        image.save(
            f"{location}/{fileName if '@' in fileName else f"{fileName.split('.')[0]}@{index}.{fileName.split('.')[-1]}"}"
        )


class ImageParser:
    Pdfs = f"{CWD}/pdfs"
    Pptx = f"{CWD}/pptxs"
    To = f"{CWD}/input"

    def trim_file(self, name: str):
        return name.replace("_", "-").replace(" ", "-").replace(",", "-")

    def parsing_pdf(self):
        pdfs = Path.iterdir(Path(ImageParser.Pdfs))
        for pdf in pdfs:
            doc: fitz.Document = fitz.open(pdf.absolute())
            for i in range(len(doc)):
                for j, img in enumerate(doc.get_page_images(i)):
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)

                    path = Path(f"{ImageParser.To}/{self.trim_file(pdf.stem)}")
                    if not path.exists():
                        path.mkdir()
                    filename = f"{path}/{pdf.stem}@{j + 1}.png"
                    pix.save(filename=filename)

    def parsing_pptx(self):
        li: Generator[Path, None, None] = Path.iterdir(Path(ImageParser.Pptx))

        for f in li:
            prs = pptx.Presentation(f)
            name = f.stem

            i = 1
            for slide in prs.slides:
                shapes: pptx.slide.SlideShapes = slide.shapes
                for shape in shapes:
                    if hasattr(shape, "image"):
                        image: pptx.ImagePart = shape.image  # type: ignore
                        image_bytes = image.blob

                        path = Path(ImageParser.To)
                        if not path.exists():
                            path.mkdir()

                        save_dir = Path(f"{path}/{self.trim_file(name)}")
                        if not save_dir.exists():
                            save_dir.mkdir()

                        path = f"{save_dir}/{name}@{i}.jpg"
                        i += 1
                        if image_bytes is not None:
                            with open(
                                path,
                                "wb",
                            ) as f:
                                f.write(image_bytes)
                        else:
                            raise Exception("image buf is empty")
